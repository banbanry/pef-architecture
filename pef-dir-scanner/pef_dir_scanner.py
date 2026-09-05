# -*- coding: utf-8 -*-
"""
PEF 目录遍历扫描归类引擎 (PEF Dir Scanner)
========================================
原理：仿 GPT token 级深度遍历 → PEF 三元拆解 → 知识库分类映射
对指定目录的文档做：
  Step 1: 遍历目录收集文档（docx/doc/pdf/md/txt/docm/pptx）
  Step 2: 提取全文（docx→python-docx, pdf→fitz, md/txt→直读, docm→zipfile正则, pptx→python-pptx）
  Step 3: token 级切块（600字/重叠120）→ PEF 三元信号打分 → 分级(S/A/B/C)
  Step 4: 文档级 PEF 三元摘要 + 知识库 10 节点分类映射 + 8层经验整理映射
  Step 5: 输出归类结果 JSON + 校对报告（对比已有分类索引）
"""
import os, re, json, math, hashlib
from collections import Counter, defaultdict

# ============ 0. 路径配置 ============
BASE = r'D:\WorkBuddy'
EXP_DIR = os.path.join(BASE, '经验整理')
B902_DIR = os.path.join(BASE, '902')
OUT_JSON = os.path.join(BASE, 'pef_scan_result.json')
OUT_REPORT = os.path.join(BASE, 'pef_scan_report.md')
DOCS_DATA = os.path.join(EXP_DIR, 'docs_data.json')

# ============ 1. PEF 三元信号词典（与 pef_loop_miner 一致） ============
P_SIGNALS = [
    '本架构', '本系统', '本方案', '本规范', 'PEF', 'P域', 'E域', 'F域', 'M层', 'M域',
    '引擎', '模型', '算子', '探针', '处理器', '模块', '平台', '框架', '底座',
    'AI', '大模型', '智能体', 'Agent', '硬件', '芯片', '单片机', 'FPGA',
    '主体', '用户', '工程师', '系统', '设备', '节点', '头雁', '从雁', '校验者',
    '仲裁者', '审查', '审计', 'CLE', 'CIC', 'GLM', 'Claude', 'GPT',
]
E_SIGNALS = [
    '变量', '输入', '输出', '参数', '状态', '数据', '阈值', '置信度', '偏差', '漂移',
    '上下文', 'token', '向量', '信号', '温度', '权重', '特征', '边界', '坐标',
    '锚', '锚点', 'π', 'pi', '哈希', '摘要', '日志', '记录', '链路', '时序',
    'E_in', 'E_out', '值', '误差', '率', '熵', '势差', '温度漂移', '温漂',
]
F_SIGNALS = [
    '结果', '结论', '裁决', '判定', '通过', '失败', 'PASS', 'FAIL', '熔断', '拒绝',
    '输出', '交付', '目标', '成果', '报告', '判决', '允许', '禁止', '闭环',
    '验证', '确认', '定案', '完成', '终止', '停机', '回退', '返回', 'ALLOW', 'DENY',
    '裁决书', '审计轨迹', '证据', '验收',
]
MECH_SIGNALS = [
    '公理', '铁则', '熔断', '审计', '锚定', '隔离', '时序', '因果', '不可伪',
    '平权', '异构', '否决', '裁决', 'MOD3', '拜占庭', '影子图', '雁阵',
    '漂移校验', '三层纪律', '第一性原理', '三元', '流水线', '状态机',
]

# ============ 2. 知识库 10 节点分类词典 ============
KB_CATEGORIES = {
    '01-算子定义与规范': ['算子', 'E033', 'E034', 'E035', 'E039', 'E040', 'E041', 'E042', 'E043',
                        'E049', 'E022', 'E056', 'E150', '探针', '检测器', 'Placeholder', 'DeadCode',
                        'BufferOverflow', 'ResourceLeak', '检测模式', '模式匹配', '正则'],
    '02-理论与设计': ['公理', 'A1', 'A2', 'A3', 'A4', 'A5', 'A6', 'A7', 'A8', 'π锚', 'pi锚',
                    'MOD3', '三元', '原语', '拓扑', '流水线', '第一性原理', '理论', '设计规范',
                    '架构论文', '架构原理', '公理体系', '时间理论', '势差', '锚定'],
    '03-测试与验证': ['测试', '验证', '故障', '缺陷', 'P0', 'P1', 'P2', 'P3', '回归', '拜占庭',
                    '测试用例', '验证报告', '审查报告', '分析报告', '实验结果', '自检', 'SELF-CHECK'],
    '04-版本与变更': ['版本', '变更', '升级', 'V2.', 'V3.', '7.6', '7.7', '8.18', 'V1.',
                    '更新日志', '版本记录', '迭代'],
    '05-知识卡片库': ['知识卡片', '卡片', '摘要', '速查', '要点', '索引'],
    '06-PEF记忆体': ['PIMEM', '记忆', '基因链', '记忆仓库', '记忆树', '记忆锚', '记忆体',
                    '遗传', '记忆检索', 'memory'],
    '07-经验整理·八层': ['方法论', '哲学', '提示词', '工程应用', '专利', '知识产权', '素材',
                        '杂项', '经验', '反思', '认知', '思维模型'],
    '08-902批次与工程案例': ['CLE', '物流', '弘信', '报关', '单证', 'AWB', '装箱单', 'SI',
                           '民爆', '雷管', '起爆', '安全管控', '交底书', '实施总结', '任务工单',
                           '部署', '落地', '工程', '案例', 'DS证据', '金丝雀'],
    '09-原始资料归档': ['原始', '归档', '素材', '草稿', '待整理', '附录'],
    'GitHub公开内容': ['公开', '导读', '同步映射', 'README', 'MIT', '白皮书', '白皮书公开'],
}

# 经验整理 8 层词典
EXP_LAYERS = {
    '01-哲学方法论层': ['方法论', '哲学', '认知', '第一性原理', '思维', '反思', '辩证', '逻辑'],
    '02-基础架构层': ['基础架构', '架构', '公理', '框架', '体系', '底座', '元架构', 'PEF架构'],
    '03-融合架构层': ['融合', '集成', '异构', '跨模型', 'CIC', '混合', '协同'],
    '04-AI基础研究层': ['AI', '大模型', '模型', '机器学习', '深度学习', '智能', '幻觉', '多模态'],
    '05-工程应用层': ['工程', '应用', '部署', '实现', '代码', '探针', 'CLE', '物流', '落地'],
    '06-提示词工程层': ['提示词', 'prompt', 'Prompt', '指令', '系统提示', 'Agent提示'],
    '07-知识产权层': ['专利', '知识产权', '交底书', '权利要求', '发明', '审核评估', 'L5'],
    '08-素材与杂项': ['素材', '杂项', '草稿', '参考', '备份', '附件', '截图'],
}

# ============ 3. 文本提取 ============
def extract_text(path):
    """按扩展名提取全文"""
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == '.docx':
            return extract_docx(path)
        elif ext == '.docm':
            return extract_docm(path)
        elif ext == '.doc':
            return extract_docx(path)  # 尝试按docx解
        elif ext == '.pdf':
            return extract_pdf(path)
        elif ext == '.md':
            with open(path, 'r', encoding='utf-8') as f:
                return f.read()
        elif ext == '.txt':
            with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        elif ext == '.pptx':
            return extract_pptx(path)
        else:
            return ''
    except Exception as e:
        return f'[EXTRACT_ERROR: {e}]'

def extract_docx(path):
    try:
        import docx
        d = docx.Document(path)
        parts = [p.text for p in d.paragraphs]
        for t in d.tables:
            for row in t.rows:
                parts.append(' | '.join(c.text for c in row.cells))
        return '\n'.join(parts)
    except ImportError:
        # 兜底：zipfile 读 document.xml
        import zipfile
        from xml.etree import ElementTree as ET
        with zipfile.ZipFile(path) as z:
            xml = z.read('word/document.xml')
        root = ET.fromstring(xml)
        ns = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}
        texts = [t.text or '' for t in root.iter('{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t')]
        return '\n'.join(texts)

def extract_docm(path):
    import zipfile, re
    try:
        with zipfile.ZipFile(path) as z:
            xml = z.read('word/document.xml').decode('utf-8', errors='ignore')
        # 去标签
        text = re.sub(r'<[^>]+>', '', xml)
        return text
    except Exception:
        return ''

def extract_pdf(path):
    try:
        import fitz
        doc = fitz.open(path)
        return '\n'.join(page.get_text() for page in doc)
    except Exception:
        return ''

def extract_pptx(path):
    try:
        from pptx import Presentation
        prs = Presentation(path)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    parts.append(shape.text)
        return '\n'.join(parts)
    except Exception:
        return ''

# ============ 4. PEF 打分与分级 ============
def pef_score(text):
    p = sum(1 for w in P_SIGNALS if w.lower() in text.lower())
    e = sum(1 for w in E_SIGNALS if w.lower() in text.lower())
    f = sum(1 for w in F_SIGNALS if w.lower() in text.lower())
    m = sum(1 for w in MECH_SIGNALS if w.lower() in text.lower())
    return {'P': p, 'E': e, 'F': f, 'M': m, 'total': p + e + f + m}

def classify_block(scores):
    t = scores['total']; m = scores['M']; f = scores['F']
    if (m >= 3) or (f >= 6 and t >= 10):
        return 'S'
    if t >= 6:
        return 'A'
    if t >= 2:
        return 'B'
    return 'C'

CHUNK_SIZE, OVERLAP = 600, 120

def token_chunk(text):
    sentences = re.split(r'(?<=[。！？!?；;\n])', text)
    chunks, cur = [], ''
    for sent in sentences:
        s = sent.strip()
        if not s:
            continue
        if len(cur) + len(s) > CHUNK_SIZE and cur:
            chunks.append(cur)
            cur = cur[-OVERLAP:] if len(cur) > OVERLAP else ''
        cur += s
    if cur.strip():
        chunks.append(cur)
    return chunks

# ============ 5. 分类映射 ============
# 文件名强规则（最高优先级：文件名语义 > 内容信号）
FILENAME_RULES = [
    # (文件名正则, 知识库分类, 经验8层)
    (r'专利|交底书|知识产权|权利|发明', '07-经验整理·八层', '07-知识产权层'),
    (r'提示词|Prompt|prompt', '02-理论与设计', '06-提示词工程层'),
    (r'算子|探针|E033|E034|E035|E039|E040|E041|E042|E043|E049|E022|E056|E150', '01-算子定义与规范', '05-工程应用层'),
    (r'物流|报关|弘信|单证|AWB|装箱单', '08-902批次与工程案例', '05-工程应用层'),
    (r'民爆|雷管|起爆|安全管控|安全架构|硬件', '08-902批次与工程案例', '05-工程应用层'),
    (r'白皮书|规范|论文|公理|原理|理论|附录|V2\.|V3\.|7\.6|7\.7|8\.18', '02-理论与设计', '02-基础架构层'),
    (r'测试|验证|审查|报告|拜占庭|故障|缺陷|P0|P1|P2|P3', '03-测试与验证', '05-工程应用层'),
    (r'记忆|PIMEM|基因链', '06-PEF记忆体', '05-工程应用层'),
    (r'CIC|跨模型|融合', '02-理论与设计', '03-融合架构层'),
    (r'方法论|哲学|认知|反思', '02-理论与设计', '01-哲学方法论层'),
    (r'AI|模型|智能|幻觉|多模态|长文本', '04-AI基础研究层', '04-AI基础研究层'),
]

def map_kb_category_filename(name):
    """文件名强规则匹配"""
    for pat, kb, exp in FILENAME_RULES:
        if re.search(pat, name):
            return kb, exp
    return None, None

def map_kb_category(text):
    """映射到知识库 10 节点"""
    scores = {}
    for cat, kws in KB_CATEGORIES.items():
        s = sum(1 for w in kws if w.lower() in text.lower())
        scores[cat] = s
    if not any(scores.values()):
        return '09-原始资料归档', 0
    top = max(scores, key=scores.get)
    return top, scores[top]

def map_exp_layer(text):
    scores = {}
    for layer, kws in EXP_LAYERS.items():
        scores[layer] = sum(1 for w in kws if w.lower() in text.lower())
    if not any(scores.values()):
        return '08-素材与杂项', 0
    top = max(scores, key=scores.get)
    return top, scores[top]

# ============ 6. 主流程 ============
def scan_dir(root, prefix=''):
    """遍历目录收集文档"""
    docs = []
    for dirpath, dirnames, filenames in os.walk(root):
        for fn in filenames:
            ext = os.path.splitext(fn)[1].lower()
            if ext in ('.docx', '.doc', '.pdf', '.md', '.txt', '.docm', '.pptx'):
                full = os.path.join(dirpath, fn)
                rel = os.path.relpath(full, BASE)
                docs.append({'path': full, 'rel': rel, 'name': fn, 'ext': ext})
    return docs

def run_scan(dirs):
    all_results = []
    for d in dirs:
        print(f"扫描目录: {d['path']} ({d['label']})")
        files = scan_dir(d['path'])
        print(f"  发现文档 {len(files)} 份")
        for fi, f in enumerate(files):
            text = extract_text(f['path'])
            if not text or len(text) < 30:
                all_results.append({**f, 'status': 'EMPTY', 'chars': len(text)})
                continue
            chunks = [c for c in token_chunk(text) if len(c) >= 30]
            blocks = []
            for ci, chunk in enumerate(chunks):
                scores = pef_score(chunk)
                blocks.append({'ci': ci, 'text': chunk, 'scores': scores,
                               'level': classify_block(scores),
                               'hash': hashlib.md5(chunk.encode()).hexdigest()[:8]})
            lv = Counter(b['level'] for b in blocks)
            sa = [b for b in blocks if b['level'] in ('S', 'A')]
            p_sum = sum(b['scores']['P'] for b in sa)
            e_sum = sum(b['scores']['E'] for b in sa)
            f_sum = sum(b['scores']['F'] for b in sa)
            m_sum = sum(b['scores']['M'] for b in sa)
            # 文档级置信度
            sa_ratio = len(sa) / len(blocks) if blocks else 0
            doc_level = ('S' if any(b['level'] == 'S' for b in blocks) else
                         'A' if sa else 'B' if blocks else 'C')
            # 分类：文件名强规则优先，内容信号兜底
            kb_from_name, exp_from_name = map_kb_category_filename(f['name'])
            if kb_from_name:
                kb_cat, kb_score = kb_from_name, 99  # 文件名命中=高置信
                exp_layer, exp_score = exp_from_name, 99
            else:
                kb_cat, kb_score = map_kb_category(text)
                exp_layer, exp_score = map_exp_layer(text)
            # 核心句
            core_sents = []
            for b in sorted(sa, key=lambda x: -x['scores']['total'])[:4]:
                sents = re.split(r'(?<=[。！？!?])', b['text'])
                best = max(sents, key=lambda s: pef_score(s)['total'], default='')
                if best and len(best) > 20:
                    core_sents.append(best.strip()[:150])
            # 主题词
            toks = Counter()
            for b in sa[:8]:
                for w in MECH_SIGNALS + list(KB_CATEGORIES.keys()):
                    if w.lower() in b['text'].lower():
                        toks[w] += 1
            top_kw = [k for k, v in toks.most_common(8) if v > 0][:8]
            result = {
                'name': f['name'], 'rel': f['rel'], 'ext': f['ext'],
                'chars': len(text), 'n_chunks': len(blocks),
                'level': doc_level, 'levels': dict(lv),
                'pef': {'P': p_sum, 'E': e_sum, 'F': f_sum, 'M': m_sum},
                'dom_pef': max(['P', 'E', 'F'], key=lambda k: {'P': p_sum, 'E': e_sum, 'F': f_sum}[k]),
                'kb_category': kb_cat, 'kb_score': kb_score,
                'exp_layer': exp_layer, 'exp_score': exp_score,
                'core': core_sents[:3], 'keywords': top_kw,
                'status': 'OK',
            }
            all_results.append(result)
            if fi % 20 == 0:
                print(f"  [{fi+1}/{len(files)}] {f['name'][:30]} → {kb_cat} ({kb_score})")
    return all_results

def write_report(results, docs_data=None):
    lines = []
    lines.append('# PEF 目录遍历扫描归类报告\n')
    lines.append(f"生成时间: 2026-09-05\n")
    lines.append(f"扫描文档: {len(results)} 份\n")
    ok = [r for r in results if r['status'] == 'OK']
    lines.append(f"有效文档: {len(ok)} 份 | 空/无法提取: {len(results) - len(ok)} 份\n")
    # 总体统计
    lines.append('## 一、知识库分类分布\n')
    cat_cnt = Counter(r['kb_category'] for r in ok)
    lines.append('| 分类 | 数量 |')
    lines.append('|---|---|')
    for cat, n in cat_cnt.most_common():
        lines.append(f'| {cat} | {n} |')
    lines.append('')
    # PEF 三元分布
    lines.append('## 二、PEF 三元主导分布\n')
    dom_cnt = Counter(r['dom_pef'] for r in ok)
    lines.append('| 主导 | 含义 | 数量 |')
    lines.append('|---|---|---|')
    lines.append(f"| P | 主体/系统/架构叙事 | {dom_cnt.get('P', 0)} |")
    lines.append(f"| E | 变量/状态/数据/锚 | {dom_cnt.get('E', 0)} |")
    lines.append(f"| F | 结果/裁决/验证/交付 | {dom_cnt.get('F', 0)} |")
    lines.append('')
    # 文档级分级
    lines.append('## 三、文档级深度分级\n')
    lv_cnt = Counter(r['level'] for r in ok)
    lines.append(f"- S 级（强机制/核心结论）: {lv_cnt.get('S', 0)}")
    lines.append(f"- A 级（机制描述）: {lv_cnt.get('A', 0)}")
    lines.append(f"- B 级（过程例子）: {lv_cnt.get('B', 0)}")
    lines.append(f"- C 级（弱信号）: {lv_cnt.get('C', 0)}")
    lines.append('')
    # 逐份清单
    lines.append('## 四、逐份归类清单\n')
    lines.append('| 文件 | 分类 | 层级 | PEF主导 | 等级 | 字符数 |')
    lines.append('|---|---|---|---|---|---|')
    for r in sorted(ok, key=lambda x: x['rel']):
        lines.append(f"| {r['name'][:40]} | {r['kb_category']} | {r['exp_layer'][:8]} | {r['dom_pef']} | {r['level']} | {r['chars']} |")
    lines.append('')
    # 校对部分（如果提供 docs_data）
    if docs_data:
        lines.append('## 五、与现有分类校对\n')
        mismatch = 0
        for r in ok:
            # 找 docs_data 中同名
            dd = next((x for x in docs_data if x.get('name') == r['name']), None)
            if dd:
                cur_cat = dd.get('cat1', '')
                # 简化：8层映射比较
                r_layer = r['exp_layer']
                if cur_cat and r_layer[:2] != cur_cat[:2]:
                    mismatch += 1
                    lines.append(f"- **{r['name'][:40]}**\n  - 现有: {cur_cat} → 扫描: {r_layer}\n")
        lines.append(f'\n共发现 {mismatch} 处潜在分类不一致（以扫描结果为准，需人工复核）。\n')
    return '\n'.join(lines)

if __name__ == '__main__':
    print("===== PEF 目录遍历扫描归类引擎 =====")
    dirs = [
        {'path': EXP_DIR, 'label': '经验整理'},
        {'path': B902_DIR, 'label': '902批次'},
    ]
    results = run_scan(dirs)
    # 保存 JSON
    with open(OUT_JSON, 'w', encoding='utf-8') as f:
        json.dump(results, f, ensure_ascii=False, indent=1)
    print(f"\n✅ 扫描完成，结果: {OUT_JSON}")
    # 校对报告
    try:
        with open(DOCS_DATA, 'r', encoding='utf-8') as f:
            docs_data = json.load(f)
    except Exception:
        docs_data = None
    report = write_report(results, docs_data)
    with open(OUT_REPORT, 'w', encoding='utf-8') as f:
        f.write(report)
    print(f"✅ 报告: {OUT_REPORT}")
